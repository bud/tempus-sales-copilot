"""Generate docs/slides.pptx — Tempus Sales Copilot pitch deck (8 slides).

Run:
    pip install python-pptx
    python create_slides.py
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Brand colors
# ---------------------------------------------------------------------------
NAVY    = RGBColor(0x0D, 0x2B, 0x4E)   # Deep navy — primary background/headers
BLUE    = RGBColor(0x1D, 0x4E, 0xD8)   # Tempus blue — accents, highlights
TEAL    = RGBColor(0x05, 0x96, 0x69)   # Success / positive
LIGHT   = RGBColor(0xEF, 0xF6, 0xFF)   # Light blue background
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x6B, 0x72, 0x80)
LGRAY   = RGBColor(0xF3, 0xF4, 0xF6)
DARK    = RGBColor(0x11, 0x18, 0x27)
AMBER   = RGBColor(0xD9, 0x77, 0x06)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def rect(slide, left, top, width, height, fill_rgb: RGBColor | None = None, line_rgb: RGBColor | None = None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(0.75)
    else:
        line.fill.background()
    return shape


def add_text(
    slide,
    text: str,
    left, top, width, height,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = DARK,
    align=PP_ALIGN.LEFT,
    italic: bool = False,
    wrap: bool = True,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_para(tf, text: str, font_size: int = 16, bold: bool = False,
             color: RGBColor = DARK, indent: int = 0, bullet: str = ""):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = indent
    run = p.add_run()
    run.text = (bullet + text) if bullet else text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def left_stripe(slide, color: RGBColor = BLUE, width: float = 0.07):
    rect(slide, Inches(0), Inches(0), Inches(width), SLIDE_H, fill_rgb=color)


def slide_number(slide, n: int):
    add_text(
        slide, str(n),
        SLIDE_W - Inches(0.5), SLIDE_H - Inches(0.35), Inches(0.35), Inches(0.3),
        font_size=9, color=GRAY, align=PP_ALIGN.RIGHT,
    )


def footer_bar(slide, text: str = "Tempus Sales Copilot  ·  Confidential"):
    rect(slide, Inches(0), SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3), fill_rgb=NAVY)
    add_text(
        slide, text,
        Inches(0.2), SLIDE_H - Inches(0.3), SLIDE_W - Inches(0.4), Inches(0.3),
        font_size=8, color=WHITE, align=PP_ALIGN.LEFT,
    )


def bullet_box(slide, items: list[str], left, top, width, height,
               font_size: int = 16, color: RGBColor = DARK, accent: RGBColor = BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet dot
        run0 = p.add_run()
        run0.text = "• "
        run0.font.size = Pt(font_size)
        run0.font.color.rgb = accent
        run0.font.bold = True
        # text
        run1 = p.add_run()
        run1.text = item
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = color
        p.space_after = Pt(6)


def stat_card(slide, value: str, label: str, left, top, width=Inches(2.2), height=Inches(1.1)):
    rect(slide, left, top, width, height, fill_rgb=LIGHT)
    add_text(slide, value, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.55),
             font_size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, label, left + Inches(0.1), top + Inches(0.65), width - Inches(0.2), Inches(0.4),
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_1_title(prs: Presentation) -> None:
    """Title slide."""
    slide = blank_slide(prs)
    # Full navy background
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=NAVY)
    # Blue accent stripe left
    rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, fill_rgb=BLUE)
    # Teal accent stripe
    rect(slide, Inches(0.08), Inches(0), Inches(0.04), SLIDE_H, fill_rgb=TEAL)

    # Logo placeholder box
    rect(slide, Inches(1.2), Inches(1.8), Inches(0.7), Inches(0.7), fill_rgb=BLUE)
    add_text(slide, "🔬", Inches(1.25), Inches(1.83), Inches(0.6), Inches(0.6),
             font_size=24, align=PP_ALIGN.CENTER, color=WHITE)

    add_text(slide, "TEMPUS SALES COPILOT",
             Inches(1.2), Inches(2.65), Inches(10), Inches(0.8),
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "From data to doorstep in 5 minutes",
             Inches(1.2), Inches(3.5), Inches(10), Inches(0.55),
             font_size=22, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.LEFT, italic=True)
    add_text(slide, "Turning territory data into personalized meeting prep, automatically.",
             Inches(1.2), Inches(4.15), Inches(9), Inches(0.5),
             font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.LEFT)

    # Bottom meta
    rect(slide, Inches(1.2), Inches(5.5), Inches(0.04), Inches(0.8), fill_rgb=TEAL)
    add_text(slide, "Case Study  ·  GenAI Product Builder Role  ·  Tempus AI",
             Inches(1.35), Inches(5.55), Inches(9), Inches(0.35),
             font_size=11, color=GRAY, align=PP_ALIGN.LEFT)
    add_text(slide, "Harish Arul  ·  2026",
             Inches(1.35), Inches(5.9), Inches(9), Inches(0.3),
             font_size=10, color=GRAY, align=PP_ALIGN.LEFT)


def slide_2_problem(prs: Presentation) -> None:
    """The Problem."""
    slide = blank_slide(prs)
    left_stripe(slide)
    footer_bar(slide)
    slide_number(slide, 2)

    add_text(slide, "THE PROBLEM", Inches(0.3), Inches(0.3), Inches(4), Inches(0.3),
             font_size=10, bold=True, color=BLUE)
    add_text(slide, "Reps have the data.\nThey don't have the time.",
             Inches(0.3), Inches(0.65), Inches(6.5), Inches(1.3),
             font_size=30, bold=True, color=DARK)

    # Stats row
    stat_card(slide, "2–4 hrs", "prep per meeting", Inches(0.3), Inches(2.15))
    stat_card(slide, "3", "disconnected\ndata sources", Inches(2.7), Inches(2.15))
    stat_card(slide, "150+", "accounts per\nrep territory", Inches(5.1), Inches(2.15))
    stat_card(slide, "~12%", "pipeline close\nrate today", Inches(7.5), Inches(2.15))

    # Bullet points
    bullet_box(
        slide,
        [
            "2–4 hours of manual prep for every 15-minute meeting",
            "3 disconnected sources: Salesforce CRM, market spreadsheets, product docs",
            "Junior reps miss connections that cost deals (e.g., high-volume lung oncologist not ordering xR RNA sequencing)",
            "Result: fewer meetings, generic pitches, inconsistent rep performance",
        ],
        Inches(0.3), Inches(3.5), Inches(9.5), Inches(2.7),
        font_size=15, color=DARK,
    )

    # Right panel — quote
    rect(slide, Inches(10.1), Inches(1.0), Inches(3.0), Inches(5.5),
         fill_rgb=RGBColor(0xF0, 0xF9, 0xFF))
    add_text(slide, "\"I have 150 accounts. Which 10 should I call this week, and what exactly should I say to each one?\"",
             Inches(10.25), Inches(1.3), Inches(2.7), Inches(2.0),
             font_size=13, color=DARK, italic=True, align=PP_ALIGN.CENTER)
    add_text(slide, "— Tempus Sales Rep",
             Inches(10.25), Inches(3.4), Inches(2.7), Inches(0.35),
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Primary User",
             Inches(10.25), Inches(4.7), Inches(2.7), Inches(0.3),
             font_size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    bullet_box(slide,
               ["Field sales rep", "~150 accounts", "8 meetings/week"],
               Inches(10.25), Inches(5.1), Inches(2.7), Inches(1.0),
               font_size=11, color=GRAY)


def slide_3_insight(prs: Presentation) -> None:
    """The Insight."""
    slide = blank_slide(prs)
    left_stripe(slide, color=TEAL)
    footer_bar(slide)
    slide_number(slide, 3)

    add_text(slide, "THE INSIGHT", Inches(0.3), Inches(0.3), Inches(4), Inches(0.3),
             font_size=10, bold=True, color=TEAL)
    add_text(slide, "This isn't a search problem.\nIt's a synthesis problem.",
             Inches(0.3), Inches(0.65), Inches(7), Inches(1.3),
             font_size=30, bold=True, color=DARK)

    # Three columns
    cols = [
        ("WHO", "to call", "Ranked by revenue opportunity × fit × recency × relationship warmth", BLUE),
        ("WHAT", "to say", "Personalized to the doctor's specialty, patient mix, and known interests", TEAL),
        ("HOW", "to handle pushback", "Grounded in real Tempus performance data, not generic talking points", AMBER),
    ]
    for i, (big, small, desc, color) in enumerate(cols):
        left = Inches(0.3 + i * 4.2)
        rect(slide, left, Inches(2.2), Inches(3.9), Inches(3.8), fill_rgb=LGRAY)
        rect(slide, left, Inches(2.2), Inches(3.9), Inches(0.06), fill_rgb=color)
        add_text(slide, big, left + Inches(0.2), Inches(2.5), Inches(1.5), Inches(0.75),
                 font_size=36, bold=True, color=color, align=PP_ALIGN.LEFT)
        add_text(slide, small, left + Inches(0.2), Inches(3.25), Inches(3.5), Inches(0.4),
                 font_size=16, color=GRAY, align=PP_ALIGN.LEFT)
        add_text(slide, desc, left + Inches(0.2), Inches(3.7), Inches(3.5), Inches(1.8),
                 font_size=13, color=DARK, align=PP_ALIGN.LEFT)

    add_text(slide,
             "The Copilot doesn't replace the rep's judgment. It gives every rep the preparation of a 10-year veteran.",
             Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.5),
             font_size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


def slide_4_solution(prs: Presentation) -> None:
    """The Solution."""
    slide = blank_slide(prs)
    left_stripe(slide)
    footer_bar(slide)
    slide_number(slide, 4)

    add_text(slide, "THE SOLUTION", Inches(0.3), Inches(0.3), Inches(4), Inches(0.3),
             font_size=10, bold=True, color=BLUE)
    add_text(slide, "Three outputs. One tool.",
             Inches(0.3), Inches(0.65), Inches(9), Inches(0.7),
             font_size=32, bold=True, color=DARK)

    outputs = [
        ("01", "Priority Rankings",
         "Ranked provider list scored by:\n• Patient volume (0.35×)\n• Product-market fit (0.25×)\n• Recency of last contact (0.20×)\n• Objection resolvability (0.15×)\n• Relationship warmth (0.05×)",
         BLUE, "Deterministic code: auditable & tunable"),
        ("02", "Objection Handler",
         "Pre-drafted responses to every known concern:\n• Grounded in real Tempus data\n• Cites source document + study\n• If KB has no answer → 'No verified data available'\n• No hallucinations",
         TEAL, "RAG pipeline: retrieve + generate"),
        ("03", "Meeting Script",
         "30-second elevator pitch per doctor:\n• Names the specific Tempus product\n• References their patient population\n• Addresses their #1 interest\n• Ends with clear call to action",
         AMBER, "LLM + personalization context"),
    ]
    for i, (num, title, body, color, tech) in enumerate(outputs):
        left = Inches(0.3 + i * 4.3)
        rect(slide, left, Inches(1.6), Inches(4.0), Inches(5.0), fill_rgb=LGRAY)
        rect(slide, left, Inches(1.6), Inches(4.0), Inches(0.08), fill_rgb=color)
        add_text(slide, num, left + Inches(0.2), Inches(1.75), Inches(0.8), Inches(0.45),
                 font_size=24, bold=True, color=color, align=PP_ALIGN.LEFT)
        add_text(slide, title, left + Inches(0.2), Inches(2.2), Inches(3.6), Inches(0.5),
                 font_size=17, bold=True, color=DARK, align=PP_ALIGN.LEFT)
        add_text(slide, body, left + Inches(0.2), Inches(2.75), Inches(3.6), Inches(2.8),
                 font_size=12, color=DARK, align=PP_ALIGN.LEFT)
        rect(slide, left + Inches(0.2), Inches(5.8), Inches(3.5), Inches(0.55),
             fill_rgb=WHITE, line_rgb=color)
        add_text(slide, tech, left + Inches(0.35), Inches(5.88), Inches(3.2), Inches(0.4),
                 font_size=10, color=color, italic=True, align=PP_ALIGN.LEFT)


def slide_5_how_it_works(prs: Presentation) -> None:
    """System flow diagram."""
    slide = blank_slide(prs)
    left_stripe(slide)
    footer_bar(slide)
    slide_number(slide, 5)

    add_text(slide, "HOW IT WORKS", Inches(0.3), Inches(0.3), Inches(4), Inches(0.3),
             font_size=10, bold=True, color=BLUE)
    add_text(slide, "Input → Synthesize → Brief",
             Inches(0.3), Inches(0.65), Inches(9), Inches(0.6),
             font_size=30, bold=True, color=DARK)

    # Flow: 3 input boxes → arrow → 2 pipeline boxes → arrow → 3 output boxes
    # Row 1: Inputs
    add_text(slide, "DATA INPUTS", Inches(0.5), Inches(1.55), Inches(3.5), Inches(0.3),
             font_size=9, bold=True, color=GRAY)
    inputs = [
        ("Market Data CSV", "Provider volume,\ncancer types, vendors"),
        ("CRM Notes JSON", "Objections, interests,\nrelationship stage"),
        ("Product KB .md", "Tempus panel specs,\ncompetitive advantages"),
    ]
    for i, (title, sub) in enumerate(inputs):
        left = Inches(0.3 + i * 1.85)
        rect(slide, left, Inches(1.9), Inches(1.65), Inches(1.1), fill_rgb=LIGHT, line_rgb=BLUE)
        add_text(slide, title, left + Inches(0.1), Inches(1.97), Inches(1.45), Inches(0.35),
                 font_size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_text(slide, sub, left + Inches(0.1), Inches(2.3), Inches(1.45), Inches(0.6),
                 font_size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # Arrow down
    add_text(slide, "▼", Inches(2.5), Inches(3.1), Inches(0.5), Inches(0.4),
             font_size=20, color=BLUE, align=PP_ALIGN.CENTER)

    # Row 2: Processing
    add_text(slide, "PROCESSING", Inches(0.5), Inches(3.5), Inches(3.5), Inches(0.3),
             font_size=9, bold=True, color=GRAY)
    processing = [
        ("Scoring Engine", "Deterministic code\n0.35×vol + 0.25×fit\n+ 0.20×recency…", NAVY),
        ("RAG Pipeline", "Embed KB chunks\nRetrieve top-3\nGenerate response", TEAL),
    ]
    for i, (title, sub, color) in enumerate(processing):
        left = Inches(0.3 + i * 2.8)
        rect(slide, left, Inches(3.8), Inches(2.5), Inches(1.3), fill_rgb=color)
        add_text(slide, title, left + Inches(0.1), Inches(3.87), Inches(2.3), Inches(0.35),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, sub, left + Inches(0.1), Inches(4.25), Inches(2.3), Inches(0.75),
                 font_size=9, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

    # LLM API box
    rect(slide, Inches(0.3), Inches(5.25), Inches(5.1), Inches(0.65), fill_rgb=LGRAY, line_rgb=GRAY)
    add_text(slide, "LLM API (GPT-4o-mini or Claude Haiku)  ·  sentence-transformers (all-MiniLM-L6-v2)  ·  numpy cosine similarity",
             Inches(0.4), Inches(5.35), Inches(4.8), Inches(0.45),
             font_size=9, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Right side: outputs description
    rect(slide, Inches(6.2), Inches(1.5), Inches(6.8), Inches(5.3), fill_rgb=LGRAY)
    add_text(slide, "WHAT THE REP SEES", Inches(6.5), Inches(1.7), Inches(6.0), Inches(0.3),
             font_size=9, bold=True, color=GRAY)

    outputs_desc = [
        ("🏆  Ranked Provider List", "10 providers sorted by priority score. Filters by stage (active/warm/cold). Score formula visible and explainable.", BLUE),
        ("⚠  Objection Handler", "Every known concern pre-addressed with Tempus-specific data. Source citations inline. 'No verified data available' when KB has no answer.", AMBER),
        ("💬  Meeting Script", "30-second pitch: names the product, references the doctor's specialty, ends with CTA. Demo fallbacks ensure the app always works.", TEAL),
    ]
    for i, (title, body, color) in enumerate(outputs_desc):
        top = Inches(2.1 + i * 1.5)
        rect(slide, Inches(6.4), top, Inches(0.06), Inches(1.1), fill_rgb=color)
        add_text(slide, title, Inches(6.6), top + Inches(0.05), Inches(6.0), Inches(0.35),
                 font_size=12, bold=True, color=color)
        add_text(slide, body, Inches(6.6), top + Inches(0.4), Inches(6.0), Inches(0.7),
                 font_size=11, color=DARK)


def slide_6_demo(prs: Presentation) -> None:
    """Demo outputs."""
    slide = blank_slide(prs)
    left_stripe(slide, color=TEAL)
    footer_bar(slide)
    slide_number(slide, 6)

    add_text(slide, "PROTOTYPE DEMO", Inches(0.3), Inches(0.3), Inches(5), Inches(0.3),
             font_size=10, bold=True, color=TEAL)
    add_text(slide, "What the rep sees in the app",
             Inches(0.3), Inches(0.65), Inches(9), Inches(0.6),
             font_size=30, bold=True, color=DARK)

    # Left panel mockup: ranked list
    rect(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(5.5), fill_rgb=WHITE,
         line_rgb=RGBColor(0xE5, 0xE7, 0xEB))
    add_text(slide, "Priority Rankings · 10 providers",
             Inches(0.45), Inches(1.6), Inches(4.2), Inches(0.3),
             font_size=9, color=GRAY)

    providers = [
        ("1", "Dr. Mei-Lin Zhao", "Thoracic · 420 pts", "95", BLUE, "Active"),
        ("2", "Dr. Sarah Chen",   "Breast · 320 pts",   "88", BLUE, "Warm"),
        ("3", "Dr. Aisha Thompson","Heme · 210 pts",    "82", TEAL, "Active"),
        ("4", "Dr. Raj Patel",    "Lung · 280 pts",     "78", BLUE, "Warm"),
        ("5", "Dr. Marcus Williams","GI · 230 pts",     "72", BLUE, "Warm"),
    ]
    for i, (rank, name, meta, score, color, stage) in enumerate(providers):
        top = Inches(2.0 + i * 0.9)
        bg = RGBColor(0xEF, 0xF6, 0xFF) if i == 0 else WHITE
        rect(slide, Inches(0.4), top, Inches(4.3), Inches(0.8), fill_rgb=bg,
             line_rgb=BLUE if i == 0 else RGBColor(0xE5, 0xE7, 0xEB))
        add_text(slide, rank, Inches(0.5), top + Inches(0.15), Inches(0.3), Inches(0.5),
                 font_size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, name, Inches(0.9), top + Inches(0.05), Inches(2.5), Inches(0.35),
                 font_size=11, bold=(i == 0), color=DARK)
        add_text(slide, meta, Inches(0.9), top + Inches(0.4), Inches(2.5), Inches(0.3),
                 font_size=9, color=GRAY)
        add_text(slide, score, Inches(4.0), top + Inches(0.15), Inches(0.6), Inches(0.5),
                 font_size=16, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Right panel mockup: briefing
    rect(slide, Inches(5.1), Inches(1.5), Inches(7.9), Inches(5.5), fill_rgb=WHITE,
         line_rgb=RGBColor(0xE5, 0xE7, 0xEB))
    add_text(slide, "Provider Briefing · Dr. Mei-Lin Zhao",
             Inches(5.25), Inches(1.6), Inches(7.5), Inches(0.3),
             font_size=9, bold=True, color=BLUE)

    sections = [
        ("Why call now:", "Champion account at Northwestern. Consistent xT CDx orderer. Expansion opportunity with xF for osimertinib resistance monitoring + onboarding new colleague Dr. Park.", DARK),
        ("Best-fit product:", "xT CDx (primary)  +  xF/xF+ (expansion)  +  xR (RNA sequencing)", TEAL),
        ("Objection handler:", "No active objections. Focus on expansion pitch.", AMBER),
        ("30-sec script:", '"Dr. Zhao, your xT CDx program is exceptional. The next step is xF liquid biopsy to track T790M resistance as your osimertinib patients progress. Serial ctDNA monitoring lets you catch resistance before it\'s clinically apparent…"', BLUE),
    ]
    for i, (label, text, color) in enumerate(sections):
        top = Inches(2.05 + i * 1.2)
        add_text(slide, label, Inches(5.3), top, Inches(7.5), Inches(0.3),
                 font_size=10, bold=True, color=color)
        rect(slide, Inches(5.3), top + Inches(0.3), Inches(7.5), Inches(0.75),
             fill_rgb=LGRAY)
        add_text(slide, text, Inches(5.45), top + Inches(0.35), Inches(7.2), Inches(0.65),
                 font_size=10, color=DARK)

    add_text(slide, "[ Screenshot placeholder: run app.py to see live demo ]",
             Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.25),
             font_size=9, color=GRAY, italic=True, align=PP_ALIGN.CENTER)


def slide_7_evaluation(prs: Presentation) -> None:
    """Evaluation & Impact."""
    slide = blank_slide(prs)
    left_stripe(slide)
    footer_bar(slide)
    slide_number(slide, 7)

    add_text(slide, "EVALUATION & IMPACT", Inches(0.3), Inches(0.3), Inches(5), Inches(0.3),
             font_size=10, bold=True, color=BLUE)
    add_text(slide, "Proving it works",
             Inches(0.3), Inches(0.65), Inches(7), Inches(0.6),
             font_size=30, bold=True, color=DARK)

    # Metrics table
    metrics = [
        ("Prep time per meeting", "2–4 hours", "< 15 minutes", "Time-tracking in pilot"),
        ("Meetings per week", "6–8", "12–15", "CRM meeting logs"),
        ("Objection resolution rate", "Anecdotal", "Tracked by type", "CRM outcome tags"),
        ("Pipeline close rate", "~12%", "18%+", "Salesforce pipeline"),
        ("Output accuracy", "—", "95%+ grounded", "Medical Affairs audit"),
    ]
    headers = ["Metric", "Current", "Target", "How to Measure"]
    col_widths = [Inches(3.0), Inches(1.6), Inches(1.8), Inches(3.0)]
    col_starts = [Inches(0.3), Inches(3.35), Inches(5.0), Inches(6.85)]
    row_height = Inches(0.5)
    header_top = Inches(1.6)

    # Header row
    for j, (header, cw, cs) in enumerate(zip(headers, col_widths, col_starts)):
        rect(slide, cs, header_top, cw - Inches(0.05), row_height, fill_rgb=NAVY)
        add_text(slide, header, cs + Inches(0.1), header_top + Inches(0.1),
                 cw - Inches(0.2), row_height - Inches(0.1),
                 font_size=10, bold=True, color=WHITE)

    for i, row in enumerate(metrics):
        top = header_top + (i + 1) * row_height
        bg = LGRAY if i % 2 == 0 else WHITE
        for j, (cell, cw, cs) in enumerate(zip(row, col_widths, col_starts)):
            rect(slide, cs, top, cw - Inches(0.05), row_height, fill_rgb=bg)
            cell_color = TEAL if j == 2 else DARK
            add_text(slide, cell, cs + Inches(0.1), top + Inches(0.1),
                     cw - Inches(0.2), row_height - Inches(0.1),
                     font_size=10, color=cell_color, bold=(j == 2))

    # Testing phases
    add_text(slide, "TESTING ROADMAP", Inches(0.3), Inches(5.0), Inches(5), Inches(0.3),
             font_size=9, bold=True, color=GRAY)
    phases = [
        ("Week 1", "Desk validation", "3 senior reps review outputs. What % usable as-is?"),
        ("Weeks 2–4", "Parallel pilot", "5 reps: Copilot for half accounts, manual for other half"),
        ("Month 2+", "Outcome tracking", "Pipeline velocity: Copilot vs. non-Copilot meetings"),
    ]
    for i, (week, phase, desc) in enumerate(phases):
        left = Inches(0.3 + i * 4.3)
        rect(slide, left, Inches(5.35), Inches(4.1), Inches(1.4), fill_rgb=LIGHT, line_rgb=BLUE)
        add_text(slide, week, left + Inches(0.15), Inches(5.45), Inches(1.5), Inches(0.3),
                 font_size=9, bold=True, color=BLUE)
        add_text(slide, phase, left + Inches(0.15), Inches(5.75), Inches(3.8), Inches(0.3),
                 font_size=12, bold=True, color=DARK)
        add_text(slide, desc, left + Inches(0.15), Inches(6.05), Inches(3.8), Inches(0.6),
                 font_size=10, color=GRAY)


def slide_8_why_tempus(prs: Presentation) -> None:
    """Why This Matters for Tempus."""
    slide = blank_slide(prs)
    # Navy background
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=NAVY)
    rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, fill_rgb=BLUE)
    rect(slide, Inches(0.08), Inches(0), Inches(0.04), SLIDE_H, fill_rgb=TEAL)

    add_text(slide, "WHY THIS MATTERS FOR TEMPUS", Inches(0.3), Inches(0.35), Inches(9), Inches(0.3),
             font_size=10, bold=True, color=RGBColor(0x93, 0xC5, 0xFD))
    add_text(slide, "Every minute saved is a meeting earned.",
             Inches(0.3), Inches(0.75), Inches(11), Inches(0.8),
             font_size=34, bold=True, color=WHITE)

    # Big stats
    stat_data = [
        ("212,000+", "NGS tests in Q2 2025 alone"),
        ("26%", "oncology volume growth YoY"),
        ("+4 meetings", "per rep per week if prep drops to 15 min"),
        ("400", "extra touchpoints/week at 100 reps"),
    ]
    for i, (val, label) in enumerate(stat_data):
        left = Inches(0.3 + (i % 2) * 6.5)
        top = Inches(1.8 + (i // 2) * 1.5)
        rect(slide, left, top, Inches(6.0), Inches(1.2),
             fill_rgb=RGBColor(0x1E, 0x3A, 0x60))
        add_text(slide, val, left + Inches(0.2), top + Inches(0.1), Inches(5.6), Inches(0.6),
                 font_size=30, bold=True, color=TEAL if i > 1 else BLUE, align=PP_ALIGN.LEFT)
        add_text(slide, label, left + Inches(0.2), top + Inches(0.7), Inches(5.6), Inches(0.4),
                 font_size=12, color=RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.LEFT)

    # Bottom message
    rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(0.04), fill_rgb=BLUE)
    add_text(slide,
             "More meetings → more conversions → more patients getting precision medicine.",
             Inches(0.3), Inches(5.25), Inches(12.7), Inches(0.45),
             font_size=16, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
    add_text(slide,
             "This compounds: if each of Tempus's 100+ sales reps gains 4 extra meetings/week, "
             "that's 400 additional physician touchpoints per week, at near-zero marginal cost.",
             Inches(0.3), Inches(5.75), Inches(12.7), Inches(0.7),
             font_size=12, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)

    add_text(slide, "Built by Harish Arul  ·  Tempus AI, GenAI Product Builder Case Study  ·  2026",
             Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.3),
             font_size=9, color=GRAY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    out_path = Path("docs/slides.pptx")
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = new_prs()
    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_insight(prs)
    slide_4_solution(prs)
    slide_5_how_it_works(prs)
    slide_6_demo(prs)
    slide_7_evaluation(prs)
    slide_8_why_tempus(prs)

    prs.save(str(out_path))
    print(f"Saved: {out_path.resolve()}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
